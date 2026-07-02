"""通用 PPT 生成 Agent —— 支持多风格主题与用户自选样式。"""
import io
import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from ..services.upload_service import save_bytes

# ========== 画布尺寸 ==========

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
MARGIN = Inches(0.9)

# ========== 基础色 ==========

C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_BLACK = RGBColor(0x11, 0x18, 0x27)
C_GRAY_700 = RGBColor(0x37, 0x41, 0x51)
C_GRAY_500 = RGBColor(0x6B, 0x72, 0x80)
C_GRAY_300 = RGBColor(0xD1, 0xD5, 0xDB)
C_GRAY_100 = RGBColor(0xF3, 0xF4, 0xF6)


STYLE_PRESETS = {
    "academic_blue": {
        "id": "academic_blue",
        "name": "学术蓝",
        "description": "稳重正式，适合通用研究汇报与阶段展示。",
        "scene": "默认推荐",
        "is_default": True,
        "fonts": {
            "title": "微软雅黑",
            "body": "微软雅黑",
        },
        "colors": {
            "primary": RGBColor(0x1A, 0x56, 0xDB),
            "secondary": RGBColor(0x0E, 0x3A, 0x9A),
            "accent": RGBColor(0xE0, 0x67, 0x1C),
            "bg": C_WHITE,
            "surface": RGBColor(0xDB, 0xE8, 0xFF),
            "text": RGBColor(0x1F, 0x29, 0x3E),
            "muted": C_GRAY_500,
            "cover_bg": RGBColor(0x0B, 0x1E, 0x4A),
        },
        "layout": {
            "cover_mode": "left_band",
            "section_mode": "split_band",
            "content_mode": "accent_bar",
            "card_mode": "soft_card",
        },
    },
    "minimal_gray": {
        "id": "minimal_gray",
        "name": "极简灰",
        "description": "留白更多，信息组织简洁，适合论文式陈述。",
        "scene": "论文型汇报",
        "is_default": False,
        "fonts": {
            "title": "微软雅黑",
            "body": "微软雅黑",
        },
        "colors": {
            "primary": RGBColor(0x4B, 0x55, 0x63),
            "secondary": RGBColor(0x1F, 0x29, 0x37),
            "accent": RGBColor(0x9C, 0xA3, 0xAF),
            "bg": C_WHITE,
            "surface": RGBColor(0xF9, 0xFA, 0xFB),
            "text": RGBColor(0x11, 0x18, 0x27),
            "muted": C_GRAY_500,
            "cover_bg": RGBColor(0xF7, 0xF7, 0xF5),
        },
        "layout": {
            "cover_mode": "minimal_center",
            "section_mode": "clean_header",
            "content_mode": "minimal_line",
            "card_mode": "outline_card",
        },
    },
    "tech_dark": {
        "id": "tech_dark",
        "name": "科技深色",
        "description": "深底高对比，更偏系统展示与技术方案表达。",
        "scene": "技术系统展示",
        "is_default": False,
        "fonts": {
            "title": "微软雅黑",
            "body": "微软雅黑",
        },
        "colors": {
            "primary": RGBColor(0x00, 0xC2, 0xA8),
            "secondary": RGBColor(0x18, 0xB6, 0xF6),
            "accent": RGBColor(0xF5, 0x9E, 0x0B),
            "bg": RGBColor(0x08, 0x11, 0x1F),
            "surface": RGBColor(0x12, 0x1D, 0x30),
            "text": RGBColor(0xE5, 0xF0, 0xFF),
            "muted": RGBColor(0x94, 0xA3, 0xB8),
            "cover_bg": RGBColor(0x05, 0x0C, 0x18),
        },
        "layout": {
            "cover_mode": "dark_grid",
            "section_mode": "tech_block",
            "content_mode": "dark_panel",
            "card_mode": "dark_card",
        },
    },
    "vibrant_orange_green": {
        "id": "vibrant_orange_green",
        "name": "活力橙绿",
        "description": "配色鲜明，强调创新点、成果和应用场景。",
        "scene": "应用导向汇报",
        "is_default": False,
        "fonts": {
            "title": "微软雅黑",
            "body": "微软雅黑",
        },
        "colors": {
            "primary": RGBColor(0xF2, 0x73, 0x29),
            "secondary": RGBColor(0x0A, 0x8F, 0x56),
            "accent": RGBColor(0x0E, 0x74, 0xC8),
            "bg": RGBColor(0xFF, 0xFB, 0xF6),
            "surface": RGBColor(0xFF, 0xF1, 0xE8),
            "text": RGBColor(0x2F, 0x24, 0x1C),
            "muted": RGBColor(0x7C, 0x6B, 0x5E),
            "cover_bg": RGBColor(0xFF, 0xF4, 0xEC),
        },
        "layout": {
            "cover_mode": "shape_stack",
            "section_mode": "color_band",
            "content_mode": "soft_panel",
            "card_mode": "vivid_card",
        },
    },
}

DEFAULT_STYLE_ID = "academic_blue"
STYLE_ORDER = [
    "academic_blue",
    "minimal_gray",
    "tech_dark",
    "vibrant_orange_green",
]


def _set_font(run, *, style: dict, role: str = "body", size: int = 18, color=None, bold=False):
    """统一设置文本样式。"""
    fonts = style["fonts"]
    colors = style["colors"]
    run.font.name = fonts["title"] if role == "title" else fonts["body"]
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or colors["text"]


def _fill_slide_background(slide, color: RGBColor):
    bg = slide.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()


def _add_page_number(slide, *, style: dict, num: int, total: int):
    colors = style["colors"]
    tb = slide.shapes.add_textbox(Inches(11.5), Inches(7.0), Inches(1.3), Inches(0.3))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num} / {total}"
    p.alignment = PP_ALIGN.RIGHT
    run = p.runs[0] if p.runs else p.add_run()
    _set_font(run, style=style, size=9, color=colors["muted"])


def _add_footer_bar(slide, *, color: RGBColor):
    bar = slide.shapes.add_shape(1, Inches(0), Inches(7.34), SLIDE_WIDTH, Inches(0.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()


def _add_cover_slide(prs, *, style: dict, topic: str, num: int, total: int):
    colors = style["colors"]
    mode = style["layout"]["cover_mode"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    if mode == "left_band":
        _fill_slide_background(slide, colors["cover_bg"])
        band = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.14), SLIDE_HEIGHT)
        band.fill.solid(); band.fill.fore_color.rgb = colors["primary"]; band.line.fill.background()
        hero = slide.shapes.add_shape(9, Inches(9.4), Inches(-1.4), Inches(4.6), Inches(4.6))
        hero.fill.solid(); hero.fill.fore_color.rgb = colors["secondary"]; hero.line.fill.background()
        mark = slide.shapes.add_shape(1, Inches(1.3), Inches(3.0), Inches(2.2), Inches(0.06))
        mark.fill.solid(); mark.fill.fore_color.rgb = colors["accent"]; mark.line.fill.background()
        title_x, title_y = Inches(1.3), Inches(2.15)
        title_color = C_WHITE
        subtitle_color = RGBColor(0x93, 0xAD, 0xD7)
        footer_color = RGBColor(0x6B, 0x8A, 0xB5)
    elif mode == "minimal_center":
        _fill_slide_background(slide, colors["cover_bg"])
        line = slide.shapes.add_shape(1, Inches(3.1), Inches(1.8), Inches(7.0), Inches(0.04))
        line.fill.solid(); line.fill.fore_color.rgb = colors["primary"]; line.line.fill.background()
        bottom = slide.shapes.add_shape(1, Inches(3.1), Inches(5.35), Inches(7.0), Inches(0.04))
        bottom.fill.solid(); bottom.fill.fore_color.rgb = colors["accent"]; bottom.line.fill.background()
        title_x, title_y = Inches(1.8), Inches(2.35)
        title_color = colors["text"]
        subtitle_color = colors["muted"]
        footer_color = colors["muted"]
    elif mode == "dark_grid":
        _fill_slide_background(slide, colors["cover_bg"])
        for x in [0.8, 3.2, 5.6, 8.0, 10.4]:
            line = slide.shapes.add_shape(1, Inches(x), Inches(0), Inches(0.02), SLIDE_HEIGHT)
            line.fill.solid(); line.fill.fore_color.rgb = colors["surface"]; line.line.fill.background()
        for y in [1.0, 2.5, 4.0, 5.5]:
            line = slide.shapes.add_shape(1, Inches(0), Inches(y), SLIDE_WIDTH, Inches(0.02))
            line.fill.solid(); line.fill.fore_color.rgb = colors["surface"]; line.line.fill.background()
        badge = slide.shapes.add_shape(1, Inches(1.0), Inches(1.1), Inches(2.0), Inches(0.28))
        badge.fill.solid(); badge.fill.fore_color.rgb = colors["primary"]; badge.line.fill.background()
        title_x, title_y = Inches(1.0), Inches(2.1)
        title_color = colors["text"]
        subtitle_color = colors["secondary"]
        footer_color = colors["muted"]
    else:  # shape_stack
        _fill_slide_background(slide, colors["cover_bg"])
        block1 = slide.shapes.add_shape(1, Inches(0.9), Inches(0.9), Inches(2.2), Inches(5.2))
        block1.fill.solid(); block1.fill.fore_color.rgb = colors["primary"]; block1.line.fill.background()
        block2 = slide.shapes.add_shape(1, Inches(3.0), Inches(1.4), Inches(1.1), Inches(4.2))
        block2.fill.solid(); block2.fill.fore_color.rgb = colors["secondary"]; block2.line.fill.background()
        dot = slide.shapes.add_shape(9, Inches(10.0), Inches(0.8), Inches(1.7), Inches(1.7))
        dot.fill.solid(); dot.fill.fore_color.rgb = colors["accent"]; dot.line.fill.background()
        title_x, title_y = Inches(4.6), Inches(2.15)
        title_color = colors["text"]
        subtitle_color = colors["secondary"]
        footer_color = colors["muted"]

    tb = slide.shapes.add_textbox(title_x, title_y, Inches(7.5), Inches(2.2))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = topic
    run = p.runs[0] if p.runs else p.add_run()
    _set_font(run, style=style, role="title", size=36, color=title_color, bold=True)
    if mode == "minimal_center":
        p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "研究汇报演示"
    if mode == "minimal_center":
        p2.alignment = PP_ALIGN.CENTER
    run2 = p2.runs[0] if p2.runs else p2.add_run()
    _set_font(run2, style=style, size=21, color=subtitle_color)

    tb2 = slide.shapes.add_textbox(Inches(1.0), Inches(5.85), Inches(10.8), Inches(0.5))
    tf2 = tb2.text_frame
    p3 = tf2.paragraphs[0]
    p3.text = "文献驱动型研究生科研 Agent  ·  通用研究汇报"
    if mode == "minimal_center":
        p3.alignment = PP_ALIGN.CENTER
    run3 = p3.runs[0] if p3.runs else p3.add_run()
    _set_font(run3, style=style, size=13, color=footer_color)

    _add_page_number(slide, style=style, num=num, total=total)


def _add_section_slide(prs, *, style: dict, title: str, color: RGBColor, num: int, total: int):
    colors = style["colors"]
    mode = style["layout"]["section_mode"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    if mode == "split_band":
        _fill_slide_background(slide, C_WHITE)
        left = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(4.5), SLIDE_HEIGHT)
        left.fill.solid(); left.fill.fore_color.rgb = color; left.line.fill.background()
        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(2.45), Inches(3.2), Inches(2.0))
        tf = title_box.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = title
        run = p.runs[0] if p.runs else p.add_run()
        _set_font(run, style=style, role="title", size=34, color=C_WHITE, bold=True)
        part_box = slide.shapes.add_textbox(Inches(5.5), Inches(2.8), Inches(6.5), Inches(1.2))
        tf2 = part_box.text_frame
        p2 = tf2.paragraphs[0]; p2.text = f"PART  {num}"
        run2 = p2.runs[0] if p2.runs else p2.add_run()
        _set_font(run2, style=style, role="title", size=44, color=color, bold=True)
    elif mode == "clean_header":
        _fill_slide_background(slide, colors["bg"])
        top = slide.shapes.add_shape(1, Inches(0.8), Inches(1.1), Inches(11.6), Inches(0.05))
        top.fill.solid(); top.fill.fore_color.rgb = color; top.line.fill.background()
        title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.35), Inches(11.0), Inches(1.0))
        tf = title_box.text_frame
        p = tf.paragraphs[0]; p.text = title; p.alignment = PP_ALIGN.CENTER
        run = p.runs[0] if p.runs else p.add_run()
        _set_font(run, style=style, role="title", size=30, color=colors["text"], bold=True)
        part_box = slide.shapes.add_textbox(Inches(4.8), Inches(4.2), Inches(3.7), Inches(0.8))
        tf2 = part_box.text_frame
        p2 = tf2.paragraphs[0]; p2.text = f"PART {num:02d}"; p2.alignment = PP_ALIGN.CENTER
        run2 = p2.runs[0] if p2.runs else p2.add_run()
        _set_font(run2, style=style, size=18, color=colors["muted"], bold=True)
    elif mode == "tech_block":
        _fill_slide_background(slide, colors["bg"])
        block = slide.shapes.add_shape(1, Inches(0.7), Inches(1.1), Inches(11.8), Inches(5.1))
        block.fill.solid(); block.fill.fore_color.rgb = colors["surface"]; block.line.fill.background()
        stripe = slide.shapes.add_shape(1, Inches(0.7), Inches(1.1), Inches(11.8), Inches(0.12))
        stripe.fill.solid(); stripe.fill.fore_color.rgb = color; stripe.line.fill.background()
        title_box = slide.shapes.add_textbox(Inches(1.2), Inches(2.1), Inches(10.5), Inches(1.2))
        tf = title_box.text_frame
        p = tf.paragraphs[0]; p.text = title
        run = p.runs[0] if p.runs else p.add_run()
        _set_font(run, style=style, role="title", size=32, color=colors["text"], bold=True)
        part_box = slide.shapes.add_textbox(Inches(1.2), Inches(4.0), Inches(5.5), Inches(1.2))
        tf2 = part_box.text_frame
        p2 = tf2.paragraphs[0]; p2.text = f"MODULE {num}"
        run2 = p2.runs[0] if p2.runs else p2.add_run()
        _set_font(run2, style=style, size=26, color=colors["secondary"], bold=True)
    else:  # color_band
        _fill_slide_background(slide, colors["bg"])
        band = slide.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_WIDTH, Inches(2.35))
        band.fill.solid(); band.fill.fore_color.rgb = color; band.line.fill.background()
        accent = slide.shapes.add_shape(1, Inches(0), Inches(2.35), SLIDE_WIDTH, Inches(0.16))
        accent.fill.solid(); accent.fill.fore_color.rgb = colors["secondary"]; accent.line.fill.background()
        title_box = slide.shapes.add_textbox(Inches(0.95), Inches(0.9), Inches(9.5), Inches(0.95))
        tf = title_box.text_frame
        p = tf.paragraphs[0]; p.text = title
        run = p.runs[0] if p.runs else p.add_run()
        _set_font(run, style=style, role="title", size=34, color=C_WHITE, bold=True)
        part_box = slide.shapes.add_textbox(Inches(9.8), Inches(0.92), Inches(2.0), Inches(0.9))
        tf2 = part_box.text_frame
        p2 = tf2.paragraphs[0]; p2.text = f"{num:02d}"; p2.alignment = PP_ALIGN.RIGHT
        run2 = p2.runs[0] if p2.runs else p2.add_run()
        _set_font(run2, style=style, role="title", size=38, color=C_WHITE, bold=True)

    _add_page_number(slide, style=style, num=num, total=total)


def _add_content_slide(prs, *, style: dict, title: str, items: list[str], color: RGBColor, num: int, total: int):
    colors = style["colors"]
    mode = style["layout"]["content_mode"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg_color = colors["surface"] if mode == "dark_panel" else colors["bg"]
    _fill_slide_background(slide, bg_color)

    if mode == "accent_bar":
        bar = slide.shapes.add_shape(1, MARGIN, Inches(0.95), Inches(0.09), Inches(0.75))
        bar.fill.solid(); bar.fill.fore_color.rgb = color; bar.line.fill.background()
        title_x, title_y = MARGIN + Inches(0.28), Inches(0.88)
        body_x, body_y, body_w = MARGIN + Inches(0.28), Inches(2.0), Inches(11.0)
    elif mode == "minimal_line":
        line = slide.shapes.add_shape(1, MARGIN, Inches(1.25), Inches(3.3), Inches(0.03))
        line.fill.solid(); line.fill.fore_color.rgb = color; line.line.fill.background()
        title_x, title_y = MARGIN, Inches(0.72)
        body_x, body_y, body_w = MARGIN, Inches(2.05), Inches(11.4)
    elif mode == "dark_panel":
        panel = slide.shapes.add_shape(1, Inches(0.72), Inches(1.32), Inches(11.85), Inches(5.45))
        panel.fill.solid(); panel.fill.fore_color.rgb = colors["surface"]; panel.line.fill.background()
        stripe = slide.shapes.add_shape(1, Inches(0.72), Inches(1.32), Inches(11.85), Inches(0.08))
        stripe.fill.solid(); stripe.fill.fore_color.rgb = color; stripe.line.fill.background()
        title_x, title_y = Inches(1.05), Inches(1.62)
        body_x, body_y, body_w = Inches(1.05), Inches(2.45), Inches(10.95)
    else:  # soft_panel
        panel = slide.shapes.add_shape(1, Inches(0.78), Inches(1.18), Inches(11.75), Inches(5.58))
        panel.fill.solid(); panel.fill.fore_color.rgb = colors["surface"]; panel.line.fill.background()
        chip = slide.shapes.add_shape(1, Inches(0.78), Inches(1.18), Inches(2.4), Inches(0.4))
        chip.fill.solid(); chip.fill.fore_color.rgb = color; chip.line.fill.background()
        title_x, title_y = Inches(1.08), Inches(1.75)
        body_x, body_y, body_w = Inches(1.08), Inches(2.52), Inches(10.9)

    title_box = slide.shapes.add_textbox(title_x, title_y, Inches(10.6), Inches(0.85))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    title_color = colors["text"]
    run = p.runs[0] if p.runs else p.add_run()
    _set_font(run, style=style, role="title", size=28, color=title_color, bold=True)

    body_box = slide.shapes.add_textbox(body_x, body_y, body_w, Inches(4.45))
    btf = body_box.text_frame; btf.word_wrap = True
    text_color = colors["text"]
    for i, item in enumerate(items):
        bp = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
        prefix = "●  " if item else ""
        bp.text = f"{prefix}{item}"
        runb = bp.runs[0] if bp.runs else bp.add_run()
        _set_font(runb, style=style, size=16, color=text_color)
        bp.space_after = Pt(8)

    _add_footer_bar(slide, color=color)
    _add_page_number(slide, style=style, num=num, total=total)


def _add_numbered_list_slide(prs, *, style: dict, title: str, items: list[str], color: RGBColor, num: int, total: int):
    colors = style["colors"]
    mode = style["layout"]["content_mode"]
    bg_color = colors["surface"] if mode == "dark_panel" else colors["bg"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_slide_background(slide, bg_color)

    title_box = slide.shapes.add_textbox(MARGIN, Inches(0.6), Inches(11.3), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]; p.text = title
    run = p.runs[0] if p.runs else p.add_run()
    _set_font(run, style=style, role="title", size=28, color=colors["text"], bold=True)

    divider = slide.shapes.add_shape(1, MARGIN, Inches(1.34), Inches(3.0), Inches(0.03))
    divider.fill.solid(); divider.fill.fore_color.rgb = color; divider.line.fill.background()

    top = Inches(1.78)
    step = Inches(0.67)
    for i, item in enumerate(items[:8]):
        y = top + step * i
        circle = slide.shapes.add_shape(9, MARGIN + Inches(0.08), y + Inches(0.04), Inches(0.42), Inches(0.42))
        circle.fill.solid(); circle.fill.fore_color.rgb = color; circle.line.fill.background()
        ctf = circle.text_frame
        cp = ctf.paragraphs[0]; cp.text = str(i + 1); cp.alignment = PP_ALIGN.CENTER
        crun = cp.runs[0] if cp.runs else cp.add_run()
        _set_font(crun, style=style, size=14, color=C_WHITE, bold=True)

        item_box = slide.shapes.add_textbox(MARGIN + Inches(0.72), y, Inches(10.4), Inches(0.5))
        itf = item_box.text_frame; itf.word_wrap = True
        ip = itf.paragraphs[0]; ip.text = item
        irun = ip.runs[0] if ip.runs else ip.add_run()
        _set_font(irun, style=style, size=16, color=colors["text"])

    _add_footer_bar(slide, color=color)
    _add_page_number(slide, style=style, num=num, total=total)


def _add_card_slide(prs, *, style: dict, title: str, items: list[str], color: RGBColor, num: int, total: int):
    colors = style["colors"]
    card_mode = style["layout"]["card_mode"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_slide_background(slide, colors["bg"])

    title_box = slide.shapes.add_textbox(MARGIN, Inches(0.48), Inches(11.2), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]; p.text = title
    run = p.runs[0] if p.runs else p.add_run()
    _set_font(run, style=style, role="title", size=28, color=colors["text"], bold=True)

    divider = slide.shapes.add_shape(1, MARGIN, Inches(1.16), Inches(3.2), Inches(0.03))
    divider.fill.solid(); divider.fill.fore_color.rgb = color; divider.line.fill.background()

    card_w = Inches(3.55)
    card_h = Inches(2.28)
    gap_x = Inches(0.25)
    gap_y = Inches(0.22)
    start_x = MARGIN + Inches(0.15)
    start_y = Inches(1.55)
    cols = min(3, max(2, len(items)))

    if card_mode == "soft_card":
        card_bg = [
            RGBColor(0xEB, 0xF2, 0xFF),
            RGBColor(0xE8, 0xF5, 0xEE),
            RGBColor(0xF5, 0xEE, 0xFF),
            RGBColor(0xFF, 0xF3, 0xE8),
            RGBColor(0xE8, 0xF0, 0xFF),
            RGBColor(0xFF, 0xEE, 0xF2),
        ]
        text_color = colors["text"]
    elif card_mode == "outline_card":
        card_bg = [C_WHITE] * 6
        text_color = colors["text"]
    elif card_mode == "dark_card":
        card_bg = [RGBColor(0x12, 0x1D, 0x30), RGBColor(0x11, 0x22, 0x28)] * 3
        text_color = colors["text"]
    else:
        card_bg = [
            RGBColor(0xFF, 0xEF, 0xE4),
            RGBColor(0xEC, 0xFB, 0xF3),
            RGBColor(0xE9, 0xF4, 0xFF),
            RGBColor(0xFF, 0xF6, 0xE6),
            RGBColor(0xF3, 0xEE, 0xFF),
            RGBColor(0xED, 0xFA, 0xF7),
        ]
        text_color = colors["text"]

    for i, item in enumerate(items[:6]):
        col = i % cols
        row = i // cols
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)

        card = slide.shapes.add_shape(1, x, y, card_w, card_h)
        card.fill.solid(); card.fill.fore_color.rgb = card_bg[i % len(card_bg)]
        if card_mode == "outline_card":
            card.line.color.rgb = color
            card.line.width = Pt(1.3)
        else:
            card.line.fill.background()

        top_bar = slide.shapes.add_shape(1, x, y, card_w, Inches(0.05))
        top_bar.fill.solid(); top_bar.fill.fore_color.rgb = color; top_bar.line.fill.background()

        num_box = slide.shapes.add_textbox(x + Inches(0.18), y + Inches(0.15), Inches(0.45), Inches(0.3))
        ntf = num_box.text_frame
        np = ntf.paragraphs[0]; np.text = f"0{i + 1}" if i < 9 else str(i + 1)
        nrun = np.runs[0] if np.runs else np.add_run()
        _set_font(nrun, style=style, size=11, color=color, bold=True)

        content_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.48), card_w - Inches(0.4), card_h - Inches(0.68))
        ctf = content_box.text_frame; ctf.word_wrap = True
        cp = ctf.paragraphs[0]; cp.text = item
        crun = cp.runs[0] if cp.runs else cp.add_run()
        _set_font(crun, style=style, size=13, color=text_color)

    _add_footer_bar(slide, color=color)
    _add_page_number(slide, style=style, num=num, total=total)


def _add_timeline_slide(prs, *, style: dict, timeline_items: list, color: RGBColor, num: int, total: int):
    colors = style["colors"]
    mode = style["layout"]["content_mode"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_slide_background(slide, colors["surface"] if mode == "dark_panel" else colors["bg"])

    title_box = slide.shapes.add_textbox(MARGIN, Inches(0.5), Inches(11.2), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]; p.text = "研究计划"
    run = p.runs[0] if p.runs else p.add_run()
    _set_font(run, style=style, role="title", size=28, color=colors["text"], bold=True)

    divider = slide.shapes.add_shape(1, MARGIN, Inches(1.18), Inches(3.0), Inches(0.03))
    divider.fill.solid(); divider.fill.fore_color.rgb = color; divider.line.fill.background()

    if not timeline_items:
        placeholder = slide.shapes.add_textbox(MARGIN, Inches(2.5), Inches(11.0), Inches(1.0))
        ptf = placeholder.text_frame
        pp = ptf.paragraphs[0]; pp.text = "（待补充）"
        prun = pp.runs[0] if pp.runs else pp.add_run()
        _set_font(prun, style=style, size=16, color=colors["muted"])
        _add_page_number(slide, style=style, num=num, total=total)
        return

    line_y = Inches(3.18)
    line_x_start = MARGIN + Inches(0.55)
    line_x_end = SLIDE_WIDTH - MARGIN - Inches(0.55)
    line = slide.shapes.add_shape(1, line_x_start, line_y, line_x_end - line_x_start, Inches(0.04))
    line.fill.solid(); line.fill.fore_color.rgb = color; line.line.fill.background()

    n = len(timeline_items[:8])
    avail = line_x_end - line_x_start
    for i, item in enumerate(timeline_items[:8]):
        x = line_x_start + int(avail * (i / max(n - 1, 1))) if n > 1 else line_x_start + int(avail / 2)
        outer = slide.shapes.add_shape(9, x - Inches(0.18), line_y - Inches(0.16), Inches(0.4), Inches(0.4))
        outer.fill.solid(); outer.fill.fore_color.rgb = C_WHITE
        outer.line.color.rgb = color; outer.line.width = Pt(3)
        inner = slide.shapes.add_shape(9, x - Inches(0.08), line_y - Inches(0.06), Inches(0.2), Inches(0.2))
        inner.fill.solid(); inner.fill.fore_color.rgb = color; inner.line.fill.background()

        phase = item.get("phase", f"阶段{i+1}") if isinstance(item, dict) else str(item)
        duration = item.get("duration", "") if isinstance(item, dict) else ""
        tasks = item.get("tasks", []) if isinstance(item, dict) else []

        top_box = slide.shapes.add_textbox(x - Inches(1.2), line_y - Inches(1.32), Inches(2.5), Inches(0.75))
        ttf = top_box.text_frame; ttf.word_wrap = True
        tp = ttf.paragraphs[0]; tp.text = phase; tp.alignment = PP_ALIGN.CENTER
        trun = tp.runs[0] if tp.runs else tp.add_run()
        _set_font(trun, style=style, size=13, color=colors["text"], bold=True)
        if duration:
            tp2 = ttf.add_paragraph(); tp2.text = duration; tp2.alignment = PP_ALIGN.CENTER
            trun2 = tp2.runs[0] if tp2.runs else tp2.add_run()
            _set_font(trun2, style=style, size=10, color=colors["muted"])

        bottom_box = slide.shapes.add_textbox(x - Inches(1.35), line_y + Inches(0.48), Inches(2.7), Inches(2.0))
        btf = bottom_box.text_frame; btf.word_wrap = True
        for j, task in enumerate(tasks[:3]):
            bp = btf.paragraphs[0] if j == 0 else btf.add_paragraph()
            bp.text = f"• {task}"
            brun = bp.runs[0] if bp.runs else bp.add_run()
            _set_font(brun, style=style, size=10, color=colors["text"])
            bp.space_after = Pt(3)

    _add_page_number(slide, style=style, num=num, total=total)


def _add_ending_slide(prs, *, style: dict, num: int, total: int):
    colors = style["colors"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_slide_background(slide, colors["cover_bg"])

    left_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.14), SLIDE_HEIGHT)
    left_bar.fill.solid(); left_bar.fill.fore_color.rgb = colors["primary"]; left_bar.line.fill.background()

    title_box = slide.shapes.add_textbox(MARGIN, Inches(2.4), Inches(11.2), Inches(1.3))
    tf = title_box.text_frame
    p = tf.paragraphs[0]; p.text = "致  谢"; p.alignment = PP_ALIGN.CENTER
    run = p.runs[0] if p.runs else p.add_run()
    _set_font(run, style=style, role="title", size=46, color=colors["text"] if style["id"] == "minimal_gray" else C_WHITE, bold=True)

    sub_box = slide.shapes.add_textbox(MARGIN, Inches(4.1), Inches(11.2), Inches(0.8))
    stf = sub_box.text_frame
    sp = stf.paragraphs[0]; sp.text = "感谢各位专家、老师批评指正！"; sp.alignment = PP_ALIGN.CENTER
    srun = sp.runs[0] if sp.runs else sp.add_run()
    sub_color = colors["muted"] if style["id"] == "minimal_gray" else colors["secondary"]
    _set_font(srun, style=style, size=20, color=sub_color)

    line = slide.shapes.add_shape(1, Inches(5.0), Inches(3.65), Inches(3.3), Inches(0.03))
    line.fill.solid(); line.fill.fore_color.rgb = colors["primary"]; line.line.fill.background()

    _add_page_number(slide, style=style, num=num, total=total)


class ProjectPPTAgent:
    """通用 PPT 生成 Agent —— 多风格主题。"""

    def __init__(self, output_dir: str = "storage/generated"):
        self.output_dir = output_dir
        os.makedirs(output_dir, exist_ok=True)

    def list_styles(self) -> list[dict]:
        """列出可供前端选择的 PPT 风格。"""
        styles = []
        for style_id in STYLE_ORDER:
            style = STYLE_PRESETS[style_id]
            styles.append({
                "id": style["id"],
                "name": style["name"],
                "description": style["description"],
                "scene": style["scene"],
                "is_default": style["is_default"],
            })
        return styles

    def resolve_style(self, template: str | None) -> dict:
        """解析风格标识，未知值回退默认风格。"""
        style_id = (template or DEFAULT_STYLE_ID).strip()
        if style_id == "default":
            style_id = DEFAULT_STYLE_ID
        return STYLE_PRESETS.get(style_id, STYLE_PRESETS[DEFAULT_STYLE_ID])

    def generate(self, design: dict, template: str = DEFAULT_STYLE_ID) -> str:
        """生成 PPT 并保存到 MinIO（本地 fallback），返回对象 key。"""
        style = self.resolve_style(template)
        prs = Presentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT

        slides = self._build_slides(design, style=style)
        total = len(slides)

        for i, slide_def in enumerate(slides):
            self._render_slide(prs, slide_def, style=style, num=i + 1, total=total)

        # 保存到 BytesIO，通过 upload_service 持久化到 MinIO
        buf = io.BytesIO()
        prs.save(buf)
        pptx_bytes = buf.getvalue()

        filename = f"project_{style['id']}_{os.urandom(4).hex()}.pptx"
        object_key = f"ppt/{filename}"
        save_bytes(pptx_bytes, object_key,
                   content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation")
        return object_key

    def _render_slide(self, prs, slide_def: dict, *, style: dict, num: int, total: int):
        stype = slide_def["type"]
        if stype == "cover":
            _add_cover_slide(prs, style=style, topic=slide_def["title"], num=num, total=total)
        elif stype == "section":
            _add_section_slide(prs, style=style, title=slide_def["title"], color=slide_def["color"], num=num, total=total)
        elif stype == "content":
            _add_content_slide(prs, style=style, title=slide_def["title"], items=slide_def["items"], color=slide_def["color"], num=num, total=total)
        elif stype == "numbered":
            _add_numbered_list_slide(prs, style=style, title=slide_def["title"], items=slide_def["items"], color=slide_def["color"], num=num, total=total)
        elif stype == "cards":
            _add_card_slide(prs, style=style, title=slide_def["title"], items=slide_def["items"], color=slide_def["color"], num=num, total=total)
        elif stype == "timeline":
            _add_timeline_slide(prs, style=style, timeline_items=slide_def["timeline"], color=slide_def["color"], num=num, total=total)
        elif stype == "ending":
            _add_ending_slide(prs, style=style, num=num, total=total)

    def _build_slides(self, design: dict, *, style: dict) -> list[dict]:
        """根据设计方案构建幻灯片定义。"""
        topic = design.get("topic", "未命名课题")
        lit = design.get("literature_review", {}) or {}
        colors = style["colors"]
        section_colors = [
            colors["primary"],
            colors["secondary"],
            colors["accent"],
            colors["primary"],
        ]

        slides: list[dict] = []
        slides.append({"type": "cover", "title": topic})

        col = section_colors[0]
        slides.append({"type": "section", "title": "研究背景与意义", "color": col})
        slides.append({"type": "content", "title": "研究背景", "items": [design.get("background", "待补充")], "color": col})
        slides.append({"type": "content", "title": "研究意义", "items": [design.get("significance", "待补充")], "color": col})
        slides.append({
            "type": "content",
            "title": "国内外研究现状",
            "items": [
                "【国际研究】" + lit.get("international", "待补充"),
                "【国内研究】" + lit.get("domestic", "待补充"),
            ],
            "color": col,
        })
        if lit.get("key_references"):
            slides.append({"type": "numbered", "title": "关键参考文献", "items": lit.get("key_references", [])[:8], "color": col})
        gaps = design.get("current_gaps", [])
        if gaps:
            slides.append({"type": "numbered", "title": "当前研究不足", "items": gaps, "color": col})

        col = section_colors[1]
        slides.append({"type": "section", "title": "研究设计与方法", "color": col})
        questions = design.get("research_questions", [])
        if questions:
            slides.append({"type": "numbered", "title": "研究问题", "items": questions, "color": col})
        objectives = design.get("objectives", [])
        if objectives:
            slides.append({"type": "numbered", "title": "研究目标", "items": objectives, "color": col})
        content = design.get("content", [])
        if content:
            slides.append({"type": "cards", "title": "研究内容", "items": self._format_content_short(content), "color": col})
        methods = design.get("methods", [])
        if methods:
            slides.append({"type": "cards", "title": "研究方法", "items": methods, "color": col})
        tech_route = design.get("technical_route", [])
        if tech_route:
            slides.append({"type": "numbered", "title": "技术路线", "items": tech_route, "color": col})
        arch = design.get("system_architecture", "")
        if arch:
            slides.append({"type": "content", "title": "实验/系统设计", "items": [arch], "color": col})

        col = section_colors[2]
        slides.append({"type": "section", "title": "创新与可行性", "color": col})
        innovations = design.get("innovation_points", [])
        if innovations:
            slides.append({"type": "cards", "title": "创新点", "items": innovations, "color": col})
        feasibility = design.get("feasibility", "")
        if feasibility:
            slides.append({"type": "content", "title": "可行性分析", "items": [feasibility], "color": col})

        col = section_colors[3]
        slides.append({"type": "section", "title": "研究计划与成果", "color": col})
        timeline = design.get("timeline", [])
        if timeline:
            slides.append({"type": "timeline", "timeline": timeline, "color": col})
        outputs = design.get("expected_outputs", [])
        if outputs:
            slides.append({"type": "numbered", "title": "预期成果", "items": outputs, "color": col})

        refs = design.get("references", []) or lit.get("key_references", [])
        if refs:
            slides.append({"type": "content", "title": "参考文献", "items": refs[:12], "color": colors["primary"]})

        slides.append({"type": "ending"})
        return slides

    def _format_content_short(self, content: list) -> list[str]:
        items: list[str] = []
        for phase in content:
            if isinstance(phase, dict):
                name = phase.get("phase", "")
                tasks = phase.get("tasks", [])
                task_str = "；".join(tasks[:3])
                items.append(f"【{name}】{task_str}")
        return items if items else ["待补充"]


ppt_agent = ProjectPPTAgent(output_dir="storage/generated")
